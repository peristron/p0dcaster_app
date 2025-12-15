# p0dcaster_app_v5.py - improved version
import streamlit as st
import os
import tempfile
import json
import requests
import io
import shutil
from pathlib import Path
from datetime import datetime
import asyncio
import PyPDF2
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import yt_dlp
import ffmpeg
from openai import OpenAI
import subprocess

# --- Streamlit Page Config ---
st.set_page_config(
    page_title="PodcastLM Studio - OS Team Testing",
    page_icon="🎧",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- ffmpeg Version Check (Debug) ---
try:
    version = subprocess.check_output(["ffmpeg", "-version"]).decode()
    st.sidebar.info(f"ffmpeg found:\n{version.splitlines()[0]}")
except Exception as e:
    st.sidebar.error(f"ffmpeg not found: {e}")

# --- Session State Defaults ---
defaults = {
    "authenticated": False,
    "script_data": None,
    "source_text": "",
    "chat_history": [],
    "notebook_content": f"# Research Notebook\n**Session Started:** {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n",
    "rehearsal_audio": None
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# --- Authentication ---
def check_password():
    if st.session_state.get("password_input", "") == st.secrets.get("APP_PASSWORD", ""):
        st.session_state.authenticated = True
    else:
        st.error("Incorrect Password")

if not st.session_state.authenticated:
    st.title("Studio Login")
    st.text_input("Enter Password", type="password", key="password_input", on_change=check_password)
    st.stop()

# --- LLM Client Factory ---
def get_llm_client(model_selection, specific_model_name, openai_key, xai_key, budget_mode):
    if budget_mode or model_selection == "Model A (OpenAI)":
        if not openai_key: return None, None, "Missing OpenAI API Key"
        return OpenAI(api_key=openai_key), "gpt-4o-mini", None
    if model_selection == "Model B (xAI Grok)":
        if not xai_key: return None, None, "Missing xAI API Key"
        model_map = {
            "Grok 4.1 Fast (Recommended)": "grok-4-1-fast-reasoning",
            "Grok 4 Full": "grok-4",
            "Grok 4 Fast": "grok-4-fast-reasoning",
            "Grok Code Fast": "grok-code-fast-1"
        }
        actual = model_map.get(specific_model_name, "grok-4-1-fast-reasoning")
        return OpenAI(api_key=xai_key, base_url="https://api.x.ai/v1"), actual, None
    return None, None, "Invalid Selection"

# --- Translation Helper ---
def translate_prompt_if_needed(client, text, target_lang):
    non_english = ["Urdu", "Arabic", "Hebrew", "Hindi", "Chinese", "Japanese", "Korean", "Russian", "Turkish"]
    if any(lang in target_lang for lang in non_english):
        try:
            res = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": f"Translate exactly to {target_lang}:\n\n{text}"}]
            )
            return res.choices[0].message.content
        except Exception as e:
            st.warning(f"Translation failed: {e}")
            return text
    return text

# --- Audio Effects ---
def create_phone_effect(input_path, output_path):
    try:
        stream = ffmpeg.input(input_path)
        stream = ffmpeg.filter(stream, 'lowpass', f=3000)
        echo = ffmpeg.input(input_path).filter('adelay', delays='120|120').filter('volume', '0.7')
        stream = ffmpeg.filter([stream, echo], 'amix', inputs=2)
        ffmpeg.output(stream, output_path, acodec='mp3', audio_bitrate='192k').run(overwrite_output=True, quiet=True)
    except Exception as e:
        st.warning(f"Phone effect failed: {e}")
        shutil.copy(input_path, output_path)

# --- Final Mixing (with ffmpeg error output) ---
def mix_final_audio(tmp_dir, script_dialogue, bg_source, selected_bg_url, uploaded_bg_file, music_ramp_up, uploaded_intro, uploaded_outro):
    tmp = Path(tmp_dir)
    inputs = []
    for i, line in enumerate(script_dialogue):
        seg_path = tmp / f"{i}.mp3"
        if not seg_path.exists() or seg_path.stat().st_size == 0:
            st.error(f"Missing or empty audio segment: {seg_path}")
            continue
        if line['speaker'] == "Caller":
            phone_path = tmp / f"phone_{i}.mp3"
            create_phone_effect(str(seg_path), str(phone_path))
            inputs.append(ffmpeg.input(str(phone_path)))
        else:
            inputs.append(ffmpeg.input(str(seg_path)))
    if not inputs:
        st.error("No valid audio segments to mix.")
        return None
    if len(inputs) > 1:
        dialogue = ffmpeg.concat(*inputs, v=0, a=1, n=len(inputs))
        dialogue = dialogue.filter('apad', pad_dur=0.4)
    else:
        dialogue = inputs[0]
    dialogue = dialogue.filter('loudnorm', I=-16, LRA=11, TP=-1.5, linear_norm=True)
    # Background music
    if bg_source != "None":
        bg_path = tmp / "bg.mp3"
        if bg_source == "Presets" and selected_bg_url:
            download_file_with_headers(selected_bg_url, str(bg_path))
        elif uploaded_bg_file:
            bg_path.write_bytes(uploaded_bg_file.getvalue())
        if bg_path.exists():
            bg = ffmpeg.input(str(bg_path))
            bg = bg.filter('aloop', loop=-1, size='2**31-1')
            bg = bg.filter('volume', 0.12)
            dialogue = ffmpeg.filter([bg, dialogue], 'amix', inputs=2, duration='longest')
            dialogue = dialogue.filter('aresample', async_='1')
    if music_ramp_up and bg_source != "None":
        silence = ffmpeg.input('anullsrc=channel_layout=stereo:sample_rate=44100', f='lavfi', t=5)
        dialogue = ffmpeg.concat(silence, dialogue, v=0, a=1)
    if uploaded_intro:
        intro = ffmpeg.input(io.BytesIO(uploaded_intro.getvalue()))
        dialogue = ffmpeg.concat(intro, dialogue, v=0, a=1)
    if uploaded_outro:
        outro = ffmpeg.input(io.BytesIO(uploaded_outro.getvalue()))
        dialogue = ffmpeg.concat(dialogue, outro, v=0, a=1)
    dialogue = dialogue.filter('afade', t='out', st='end-5', d=5)
    out_path = tmp / "podcast.mp3"
    try:
        ffmpeg.output(dialogue, str(out_path), acodec='mp3', audio_bitrate='192k').run(overwrite_output=True, quiet=True)
    except ffmpeg.Error as e:
        st.warning(f"Advanced mixing failed: {e.stderr.decode(errors='ignore')}")
        # fallback...
        try:
            simple = ffmpeg.concat(*inputs, v=0, a=1, n=len(inputs))
            ffmpeg.output(simple, str(out_path), acodec='mp3', audio_bitrate='192k').run(overwrite_output=True, quiet=True)
        except Exception as e2:
            st.error(f"Simple concat also failed: {e2}")
            return None
    return out_path

# --- File Extraction ---
def extract_text_from_files(files, audio_client=None):
    text = ""
    for file in files:
        try:
            name = file.name.lower()
            if name.endswith(".pdf"):
                reader = PyPDF2.PdfReader(io.BytesIO(file.getvalue()))
                for page in reader.pages:
                    text += (page.extract_text() or "") + "\n"
            elif name.endswith(".docx"):
                doc = docx.Document(io.BytesIO(file.getvalue()))
                for para in doc.paragraphs:
                    text += para.text + "\n"
            elif name.endswith(".pptx"):
                prs = Presentation(io.BytesIO(file.getvalue()))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif name.endswith(".txt"):
                text += file.getvalue().decode("utf-8") + "\n"
            elif name.endswith((".mp3", ".wav", ".m4a", ".mp4", ".webm")):
                if audio_client:
                    with st.spinner(f"Transcribing {file.name}..."):
                        transcript = audio_client.audio.transcriptions.create(model="whisper-1", file=(file.name, file.getvalue()))
                        text += transcript.text + "\n"
                else:
                    st.warning(f"OpenAI key required for {file.name}")
            else:
                st.warning(f"Unsupported: {file.name}")
        except Exception as e:
            st.error(f"Error reading {file.name}: {e}")
    return text

def download_file_with_headers(url, save_path):
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        r = requests.get(url, headers=headers, stream=True, timeout=15)
        if r.status_code == 200:
            with open(save_path, 'wb') as f:
                for chunk in r.iter_content(8192):
                    f.write(chunk)
            return True
    except Exception as e:
        st.warning(f"Download failed: {e}")
    return False

def scrape_website(url):
    try:
        headers = {'User-Agent': 'Mozilla/5.0'}
        r = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(r.content, 'html.parser')
        for tag in soup(["script", "style", "header", "footer", "nav"]):
            tag.decompose()
        return soup.get_text()
    except Exception as e:
        st.warning(f"Scraping failed: {e}")
        return None

async def download_and_transcribe_video(url, audio_client):
    loop = asyncio.get_event_loop()
    try:
        with tempfile.TemporaryDirectory() as tmp_dir:
            ydl_opts = {
                'format': 'bestaudio/best',
                'outtmpl': os.path.join(tmp_dir, 'audio.%(ext)s'),
                'postprocessors': [{'key': 'FFmpegExtractAudio', 'preferredcodec': 'mp3'}],
                'quiet': True,
                'http_headers': {'User-Agent': 'Mozilla/5.0'}
            }
            await loop.run_in_executor(None, lambda: yt_dlp.YoutubeDL(ydl_opts).download([url]))
            audio_path = next(Path(tmp_dir).glob("audio.*"))
            with open(audio_path, "rb") as f:
                transcript = audio_client.audio.transcriptions.create(model="whisper-1", file=f)
            return transcript.text, None
    except Exception as e:
        return None, str(e)

# --- Audio Generation Helper (with debug) ---
def generate_audio_openai(client, text, voice, filename, speed=1.0):
    try:
        response = client.audio.speech.create(model="tts-1", voice=voice, input=text, speed=speed)
        response.stream_to_file(filename)
        return True
    except Exception as e:
        st.warning(f"TTS failed for voice '{voice}': {e}")
        return False

# --- Sidebar ---
with st.sidebar:
    st.title("Studio Settings")
    openai_key = st.secrets.get("OPENAI_API_KEY") or st.text_input("OpenAI API Key", type="password")
    xai_key = st.secrets.get("XAI_API_KEY") or st.text_input("xAI API Key (Optional)", type="password")
    model_choice = st.radio("Intelligence Engine", ["Model A (OpenAI)", "Model B (xAI Grok)"])
    xai_version = "Grok 4.1 Fast (Recommended)"
    if model_choice == "Model B (xAI Grok)":
        xai_version = st.selectbox("Grok Model", ["Grok 4.1 Fast (Recommended)", "Grok 4 Full", "Grok 4 Fast", "Grok Code Fast"])
    budget_mode = st.checkbox("Budget Mode (GPT-4o-mini)", help="Saves 70–90% on LLM cost")
    privacy_mode = st.toggle("Privacy Mode", value=False)
    if st.button("New Session"):
        for key in defaults:
            st.session_state[key] = defaults[key]
        st.rerun()
    st.divider()
    st.subheader("Language & Length")
    language = st.selectbox("Output Language", [
        "English (US)", "English (UK)", "Spanish", "French", "German", "Italian", "Portuguese",
        "Hindi", "Urdu", "Arabic", "Hebrew", "Russian", "Turkish", "Japanese", "Korean", "Chinese (Mandarin)"
    ])
    length_option = st.select_slider("Duration", ["Short (2 min)", "Medium (5 min)", "Long (15 min)", "Extra Long (30 min)"])
    st.subheader("Hosts")
    host1_persona = st.text_input("Host 1 Persona", "Male, curious, slightly skeptical")
    host2_persona = st.text_input("Host 2 Persona", "Female, enthusiastic expert")
    voice_style = st.selectbox("Voice Pair", ["Dynamic (Alloy & Nova)", "Calm (Onyx & Shimmer)", "Formal (Echo & Fable)"])
    voice_map = {"Dynamic (Alloy & Nova)": ("alloy", "nova"), "Calm (Onyx & Shimmer)": ("onyx", "shimmer"), "Formal (Echo & Fable)": ("echo", "fable")}
    st.divider()
    st.subheader("Music")
    bg_source = st.radio("Background Music", ["Presets", "Upload Custom", "None"], horizontal=True)
    music_ramp_up = st.checkbox("Start Music 5s Early")
    selected_bg_url = None
    uploaded_bg_file = None
    if bg_source == "Presets":
        music_choice = st.selectbox("Track", ["Lo-Fi (Study)", "Upbeat (Morning)", "Ambient (News)", "Cinematic (Deep)"])
        music_urls = {
            "Lo-Fi (Study)": "https://cdn.pixabay.com/download/audio/2022/05/27/audio_1808fbf07a.mp3",
            "Upbeat (Morning)": "https://cdn.pixabay.com/download/audio/2024/05/24/audio_95e3f5f471.mp3",
            "Ambient (News)": "https://cdn.pixabay.com/download/audio/2022/03/10/audio_c8c8a73467.mp3",
            "Cinematic (Deep)": "https://cdn.pixabay.com/download/audio/2022/03/22/audio_c2b86c77ce.mp3"
        }
        selected_bg_url = music_urls[music_choice]
    elif bg_source == "Upload Custom":
        uploaded_bg_file = st.file_uploader("Upload Loop", type=["mp3", "wav"])
    with st.expander("Intro/Outro"):
        uploaded_intro = st.file_uploader("Intro", type=["mp3", "wav"])
        uploaded_outro = st.file_uploader("Outro", type=["mp3", "wav"])
    st.divider()
    st.subheader("Live Cost Estimate, just FYI (populates AFTER script generation step)")
    if st.session_state.script_data:
        chars = sum(len(l["text"]) for l in st.session_state.script_data["dialogue"])
        tts_cost = (chars / 1_000_000) * 15
        llm_cost = 0.30 if not budget_mode and model_choice == "Model B (xAI Grok)" else 0.10
        st.metric("TTS Cost", f"${tts_cost:.3f}")
        st.metric("LLM Cost", f"${llm_cost:.2f}")
        st.success(f"Total ≈ ${tts_cost + llm_cost:.2f}")

# --- Main App ---
st.title("PodcastLM Studio - OS Team Testing")
tab1, tab2, tab3, tab4 = st.tabs(["1. Source", "2. Research Chat", "3. Script & Rehearsal", "4. Produce"])
audio_client = OpenAI(api_key=openai_key) if openai_key else None

# === TAB 1: SOURCE ===
with tab1:
    st.info("Upload content — drives both chat and podcast")
    input_type = st.radio("Input Type", ["Files", "Web URL", "Video URL", "Text"], horizontal=True)
    new_text = ""
    if input_type == "Files":
        files = st.file_uploader("Upload", accept_multiple_files=True)
        if files and st.button("Process"):
            with st.spinner("Extracting text..."):
                new_text = extract_text_from_files(files, audio_client)
    elif input_type == "Web URL":
        url = st.text_input("Article URL")
        if url and st.button("Scrape"):
            with st.spinner("Scraping..."):
                new_text = scrape_website(url) or ""
    elif input_type == "Video URL":
        vid_url = st.text_input("YouTube / Video URL")
        if vid_url and st.button("Transcribe"):
            if audio_client:
                with st.spinner("Transcribing video..."):
                    loop = asyncio.new_event_loop()
                    text, err = loop.run_until_complete(download_and_transcribe_video(vid_url, audio_client))
                    new_text = text or ""
            else:
                st.error("OpenAI key required")
    elif input_type == "Text":
        new_text = st.text_area("Paste text", height=300)
    if new_text and new_text != st.session_state.source_text:
        st.session_state.source_text = new_text
        st.session_state.chat_history = []
        st.session_state.notebook_content += f"\n---\n### New Source ({datetime.now().strftime('%H:%M')})\n\n"
        st.success("Source loaded!")

# === TAB 2: RESEARCH CHAT ===
with tab2:
    st.header("Research Chat")
    if not st.session_state.source_text:
        st.info("Please upload or enter source content in Tab 1 first.")
    else:
        # Display chat history
        for entry in st.session_state.chat_history:
            if entry["role"] == "user":
                st.markdown(f"**You:** {entry['content']}")
            else:
                st.markdown(f"**AI:** {entry['content']}")

        # Chat input
        user_question = st.text_input("Ask a question about the source:", key="research_chat_input")
        if st.button("Send", key="research_chat_send") and user_question.strip():
            # Add user question to chat history
            st.session_state.chat_history.append({"role": "user", "content": user_question})

            # Prepare LLM client
            client, model, err = get_llm_client(
                model_choice, xai_version, openai_key, xai_key, budget_mode
            )
            if err:
                st.error(err)
            else:
                # Compose prompt with source and chat history
                chat_prompt = [
                    {"role": "system", "content": "You are a helpful research assistant. Use the provided source to answer questions."},
                    {"role": "system", "content": f"Source:\n{st.session_state.source_text[:40000]}"}
                ]
                # Add previous exchanges (optional, for context)
                for entry in st.session_state.chat_history:
                    chat_prompt.append({"role": entry["role"], "content": entry["content"]})

                # Query LLM
                with st.spinner("Thinking..."):
                    try:
                        response = client.chat.completions.create(
                            model=model,
                            messages=chat_prompt,
                            max_tokens=512
                        )
                        ai_reply = response.choices[0].message.content
                        st.session_state.chat_history.append({"role": "assistant", "content": ai_reply})
                        st.markdown(f"**AI:** {ai_reply}")
                    except Exception as e:
                        st.error(f"LLM error: {e}")

        # Option to clear chat
        if st.button("Clear Chat", key="research_chat_clear"):
            st.session_state.chat_history = []

# === TAB 3: SCRIPT GENERATION ===
with tab3:
    col_dir, col_call = st.columns([1, 1])
    with col_dir:
        user_instructions = st.text_area("Director Notes", placeholder="e.g., Make it funny")
    with col_call:
        caller_prompt = st.text_area("Caller Question (optional)")
    if st.button("Generate Script", type="primary"):
        if not st.session_state.source_text:
            st.error("Load source first")
        else:
            client, model, err = get_llm_client(model_choice, xai_version, openai_key, xai_key, budget_mode)
            if err:
                st.error(err)
            else:
                with st.spinner("Writing script..."):
                    word_targets = {
                        "Short (2 min)": 800,
                        "Medium (5 min)": 2200,
                        "Long (15 min)": 6000,
                        "Extra Long (30 min)": 12000
                    }
                    target_words = word_targets[length_option]
                    length_instr = f"Write a very detailed, natural, conversational podcast script with approximately {target_words} total words ({length_option}). Use long explanations, tangents, humor, and back-and-forth dialogue. NEVER truncate lines."
                    call_in = f"Include a Caller asking: '{caller_prompt}' and hosts respond." if caller_prompt else ""
                    translated = translate_prompt_if_needed(client, user_instructions, language)
                    prompt = f"""Create a podcast script in {language}.
Host 1: {host1_persona}
Host 2: {host2_persona}
{length_instr}
Director notes: {translated}
{call_in}
Output strict JSON: {{"title": "...", "dialogue": [{{"speaker": "Host 1", "text": "..."}}, ...]}}
Source: {st.session_state.source_text[:40000]}"""
                    res = client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": prompt}],
                        response_format={"type": "json_object"}
                    )
                    st.session_state.script_data = json.loads(res.choices[0].message.content)
                    st.success("Script ready!")
                    if privacy_mode:
                        st.session_state.source_text = ""
    # Script editing and rehearsal
    if st.session_state.script_data:
        data = st.session_state.script_data
        st.subheader(data.get("title", "Untitled Podcast"))
        with st.expander("Edit Script"):
            with st.form("edit_form"):
                new_d = []
                for i, line in enumerate(data["dialogue"]):
                    c1, c2 = st.columns([1, 5])
                    speakers = ["Host 1", "Host 2"] + (["Caller"] if any(l["speaker"] == "Caller" for l in data["dialogue"]) else [])
                    speaker = c1.selectbox("Speaker", speakers, index=speakers.index(line["speaker"]) if line["speaker"] in speakers else 0, key=f"s{i}")
                    text = c2.text_area("Line", line["text"], height=80, key=f"t{i}")
                    new_d.append({"speaker": speaker, "text": text})
                if st.form_submit_button("Save"):
                    st.session_state.script_data["dialogue"] = new_d
                    st.success("Saved")
        st.subheader("Live Rehearsal")
        idx = st.selectbox("Preview line", range(len(data["dialogue"])), format_func=lambda i: f"{data['dialogue'][i]['speaker']}: {data['dialogue'][i]['text'][:60]}...")
        if st.button("Play Line"):
            if openai_key:
                line = data["dialogue"][idx]
                voice = voice_map[voice_style][0 if line["speaker"] == "Host 1" else 1]
                if line["speaker"] == "Caller":
                    voice = "fable"
                with tempfile.NamedTemporaryFile(suffix=".mp3") as tmp:
                    generate_audio_openai(audio_client, line["text"], voice, tmp.name)
                    st.audio(tmp.name)
            else:
                st.error("OpenAI key needed")

# === TAB 4: PRODUCTION (with debug) ===
with tab4:
    if st.session_state.script_data and st.button("Produce Final Podcast", type="primary"):
        if not openai_key:
            st.error("OpenAI key required")
            st.stop()
        progress = st.progress(0)
        status = st.empty()
        audio_client = OpenAI(api_key=openai_key)
        m_voice, f_voice = voice_map[voice_style]
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp = Path(tmp_dir)
            script = st.session_state.script_data["dialogue"]

            # --- DEBUG: Show script data ---
            st.write("Script data for debugging:", st.session_state.script_data)

            # --- Generate audio segments with debug checks ---
            for i, line in enumerate(script):
                status.text(f"Voicing line {i+1}/{len(script)}")
                voice = m_voice if line["speaker"] == "Host 1" else f_voice
                if line["speaker"] == "Caller":
                    voice = "fable"
                f_path = tmp / f"{i}.mp3"
                success = generate_audio_openai(audio_client, line["text"], voice, str(f_path))
                # --- DEBUG: Check audio file existence and size ---
                if not success or not f_path.exists() or f_path.stat().st_size == 0:
                    st.error(f"Audio generation failed for line {i}: {line['text'][:60]}")
                progress.progress((i + 1) / len(script))

            status.text("Mixing podcast...")
            out_path = mix_final_audio(tmp_dir, script, bg_source, selected_bg_url, uploaded_bg_file, music_ramp_up, uploaded_intro, uploaded_outro)
            if out_path and out_path.exists():
                with open(out_path, "rb") as f:
                    audio_bytes = f.read()
                status.success("Complete!")
                st.audio(audio_bytes, format="audio/mp3")
                st.download_button("Download Podcast", audio_bytes, "podcast.mp3", "audio/mp3")
            else:
                st.error("Podcast production failed. See errors above for details.")

